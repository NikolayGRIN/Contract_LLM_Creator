import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

# -----------------------
# Config
# -----------------------
KEEP_METHODS = ["bm25", "emb_llama", "hybrid_llama"]
METHOD_LABELS = {
    "bm25": "BM25",
    "emb_llama": "Embeddings (bge-m3)",
    "hybrid_llama": "Hybrid (BM25+Emb)",
}
METRICS = ["recall@5", "recall@10", "precision@10", "mrr@10"]

# -----------------------
# Locate latest CSV
# -----------------------
debug_dir = Path("debug")
csv_files = sorted(debug_dir.glob("retrieval_metrics_long_*.csv"))
if not csv_files:
    raise FileNotFoundError("Не найден debug/retrieval_metrics_long_*.csv")

csv_path = csv_files[-1]
print("Using:", csv_path)

df = pd.read_csv(csv_path)

# -----------------------
# Filter methods
# -----------------------
df = df[df["method"].isin(KEEP_METHODS)].copy()
if df.empty:
    raise ValueError(f"После фильтра KEEP_METHODS={KEEP_METHODS} данных не осталось. "
                     f"Проверь названия методов в CSV: {sorted(pd.read_csv(csv_path)['method'].unique())}")

# -----------------------
# Pivot table (method x metric)
# -----------------------
pivot = df.pivot(index="method", columns="metric", values="value")

# Ensure metric columns exist
missing_metrics = [m for m in METRICS if m not in pivot.columns]
if missing_metrics:
    raise ValueError(f"В CSV отсутствуют метрики: {missing_metrics}. Доступно: {list(pivot.columns)}")

# Reorder methods and rename labels
pivot = pivot.reindex(KEEP_METHODS)
pivot.index = [METHOD_LABELS.get(m, m) for m in pivot.index]

# Reorder metrics
pivot = pivot[METRICS]

print("\n=== TABLE (BM25 vs Embeddings vs Hybrid+Emb) ===")
print(pivot.round(3).to_string())

# Save table to Excel (optional but useful)
xlsx_path = debug_dir / "retrieval_metrics_table_3methods.xlsx"
pivot.round(6).to_excel(xlsx_path)
print("\nSaved table:", xlsx_path)

# -----------------------
# Academic combined chart (grouped bars)
# -----------------------
methods = list(pivot.index)
x = np.arange(len(methods))

width = 0.18
offsets = [-1.5 * width, -0.5 * width, 0.5 * width, 1.5 * width]

plt.figure(figsize=(11.5, 5.8))
ax = plt.gca()

bar_containers = []
for off, metric in zip(offsets, METRICS):
    vals = pivot[metric].values
    bc = ax.bar(x + off, vals, width, label=metric)
    bar_containers.append((bc, vals))

ax.set_ylim(0, 1.0)
ax.set_ylabel("Score", fontsize=12)
ax.set_title("Сравнение Retrieval: BM25 vs Embeddings vs Hybrid", fontsize=14)

ax.set_xticks(x)
ax.set_xticklabels(methods, rotation=15, ha="right", fontsize=11)

# Light grid (academic)
ax.grid(axis="y", linestyle="--", linewidth=0.6, alpha=0.6)
ax.set_axisbelow(True)

# Legend on top, compact
ax.legend(
    ncols=2,
    frameon=False,
    loc="upper right",
    fontsize=11
)

# Value labels above bars
def add_labels(container, values):
    for rect, v in zip(container, values):
        h = rect.get_height()
        ax.text(
            rect.get_x() + rect.get_width() / 2,
            h + 0.015,
            f"{float(v):.3f}",
            ha="center",
            va="bottom",
            fontsize=9,
            clip_on=False,
        )

for bc, vals in bar_containers:
    add_labels(bc, vals)

plt.tight_layout()

img_path = debug_dir / "retrieval_summary_3methods_academic.png"
plt.savefig(img_path, dpi=220)
plt.close()
print("Saved image:", img_path)

# -----------------------
# PPTX slide
# -----------------------
prs = Presentation()
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.3), Inches(0.55))
tf = title_box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = "Retrieval: BM25 vs Embeddings vs Hybrid (BM25+Emb)"
p.font.size = Pt(22)
p.font.bold = True

# Subtitle (try to read meta from df columns)
subtitle = ""
try:
    meta_row = df.iloc[0]
    qn = int(meta_row.get("queries_used", 0))
    cand = int(meta_row.get("candidates", 0))
    K = int(meta_row.get("K", 10))
    alpha = float(meta_row.get("hybrid_alpha", 0.5))
    subtitle = f"Queries: {qn} | Candidates: {cand} | K={K} | α={alpha}"
except Exception:
    pass

if subtitle:
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9.3), Inches(0.3))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(14)

# Insert image
slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.05), width=Inches(9.2))

pptx_path = debug_dir / "retrieval_summary_slide_3methods_academic.pptx"
prs.save(pptx_path)
print("Saved PPTX:", pptx_path)
